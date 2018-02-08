from pptx import Presentation
from docx import Document


document = Document()

prs = Presentation("lecture_slides.pptx")
for slide in prs.slides:
	for shape in slide.shapes:
		try:

			document.add_paragraph(shape.text)

		except AttributeError:
			try:
				document.add_picture(shape.picture)
			except AttributeError:
				continue

document.save("notes.docx")
